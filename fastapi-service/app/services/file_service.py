import logging
from typing import List

from app.models.schemas import AnalysisResponse, PracticeQuestion
from app.utils.pdf_utils   import extract_text_from_pdf, get_pdf_metadata
from app.utils.image_utils import get_image_mime_type
from app.utils.docx_utils  import extract_text_from_docx, get_docx_metadata, get_docx_headings
from app.services.ai_service import (
    analyze_text_content,
    analyze_image_content,
    generate_practice_questions,
    check_plagiarism,
)

logger = logging.getLogger(__name__)


# ─── MAIN ORCHESTRATOR ───────────────────────────────────────────────────────

def process_file(file_bytes: bytes, file_name: str, content_type: str) -> AnalysisResponse:
    """
    Routes file to the correct pipeline based on extension / content_type.
    Supported: PDF, DOCX, PPT/PPTX, Images (jpg/png/gif/webp/bmp)
    """
    try:
        ext = file_name.lower().split(".")[-1]

        if content_type == "pdf"  or ext == "pdf":
            return _process_pdf(file_bytes, file_name)

        if content_type == "docx" or ext in ("doc", "docx"):
            return _process_docx(file_bytes, file_name)

        if content_type == "image" or ext in ("jpg", "jpeg", "png", "gif", "webp", "bmp"):
            return _process_image(file_bytes, file_name)

        if content_type == "ppt"  or ext in ("ppt", "pptx"):
            return _process_ppt(file_bytes, file_name)

        # Unknown — try PDF extraction as last resort
        logger.warning(f"Unknown type '{content_type}' for '{file_name}', trying PDF pipeline")
        return _process_pdf(file_bytes, file_name)

    except Exception as e:
        logger.error(f"process_file crashed for {file_name}: {e}")
        return AnalysisResponse(
            file_name=file_name,
            content_type=content_type,
            status="error",
            error_message="An unexpected error occurred during processing.",
            confidence=0.0
        )


# ─── PDF PIPELINE ────────────────────────────────────────────────────────────

def _process_pdf(file_bytes: bytes, file_name: str) -> AnalysisResponse:
    logger.info(f"[PDF] Processing: {file_name}")

    text     = extract_text_from_pdf(file_bytes)
    metadata = get_pdf_metadata(file_bytes)

    if not text:
        return AnalysisResponse(
            file_name=file_name, content_type="pdf", status="error",
            error_message="Could not extract text. PDF may be scanned/image-based.",
            confidence=0.0
        )

    word_count = len(text.split())

    # Run all AI features
    ai        = analyze_text_content(text, file_name)
    questions = generate_practice_questions(text, file_name, count=5)
    plagiarism = check_plagiarism(text, file_name)

    return AnalysisResponse(
        file_name=file_name,
        content_type="pdf",
        # Core
        summary=ai.get("summary"),
        tags=ai.get("tags", []),
        keywords=ai.get("keywords", []),
        difficulty=ai.get("difficulty"),
        language=ai.get("language", "English"),
        topics=ai.get("topics", []),
        suggested_year=ai.get("suggested_year"),
        confidence=ai.get("confidence", 0.8),
        # Extended
        quality_score=ai.get("quality_score"),
        quality_note=ai.get("quality_note"),
        recommendations=ai.get("recommendations", []),
        practice_questions=_to_question_models(questions),
        plagiarism_score=plagiarism.get("plagiarism_score"),
        plagiarism_note=plagiarism.get("plagiarism_note"),
        # Meta
        word_count=word_count,
        page_count=metadata.get("page_count"),
        status="success"
    )


# ─── DOCX PIPELINE ───────────────────────────────────────────────────────────

def _process_docx(file_bytes: bytes, file_name: str) -> AnalysisResponse:
    logger.info(f"[DOCX] Processing: {file_name}")

    text     = extract_text_from_docx(file_bytes)
    metadata = get_docx_metadata(file_bytes)
    headings = get_docx_headings(file_bytes)

    if not text:
        return AnalysisResponse(
            file_name=file_name, content_type="docx", status="error",
            error_message="Could not extract text from DOCX. File may be empty or corrupted.",
            confidence=0.0
        )

    word_count = len(text.split())

    # Run all AI features (same as PDF)
    ai         = analyze_text_content(text, file_name)
    questions  = generate_practice_questions(text, file_name, count=5)
    plagiarism = check_plagiarism(text, file_name)

    return AnalysisResponse(
        file_name=file_name,
        content_type="docx",
        # Core
        summary=ai.get("summary"),
        tags=ai.get("tags", []),
        keywords=ai.get("keywords", []),
        difficulty=ai.get("difficulty"),
        language=ai.get("language", "English"),
        topics=ai.get("topics", []),
        suggested_year=ai.get("suggested_year"),
        confidence=ai.get("confidence", 0.8),
        # Extended
        quality_score=ai.get("quality_score"),
        quality_note=ai.get("quality_note"),
        recommendations=ai.get("recommendations", []),
        practice_questions=_to_question_models(questions),
        plagiarism_score=plagiarism.get("plagiarism_score"),
        plagiarism_note=plagiarism.get("plagiarism_note"),
        # DOCX-specific
        docx_metadata=metadata,
        headings=headings,
        tables_count=metadata.get("tables_count"),
        images_count=metadata.get("images_count"),
        # Meta
        word_count=word_count,
        status="success"
    )


# ─── IMAGE PIPELINE ──────────────────────────────────────────────────────────

def _process_image(file_bytes: bytes, file_name: str) -> AnalysisResponse:
    logger.info(f"[IMAGE] Processing: {file_name}")

    mime_type  = get_image_mime_type(file_name)
    ai         = analyze_image_content(file_bytes, file_name, mime_type)
    ocr_text   = ai.get("ocr_text", "")
    word_count = len(ocr_text.split()) if ocr_text else 0

    # Generate questions from OCR text if substantial
    questions = []
    if ocr_text and len(ocr_text.split()) > 50:
        questions = generate_practice_questions(ocr_text, file_name, count=3)

    return AnalysisResponse(
        file_name=file_name,
        content_type="image",
        summary=ai.get("summary"),
        tags=ai.get("tags", []),
        keywords=ai.get("keywords", []),
        difficulty=ai.get("difficulty"),
        topics=ai.get("topics", []),
        ocr_text=ocr_text,
        quality_score=ai.get("quality_score"),
        confidence=ai.get("confidence", 0.8),
        word_count=word_count,
        language="English",
        practice_questions=_to_question_models(questions),
        status="success"
    )


# ─── POWERPOINT PIPELINE ─────────────────────────────────────────────────────

def _process_ppt(file_bytes: bytes, file_name: str) -> AnalysisResponse:
    logger.info(f"[PPT] Processing: {file_name}")

    try:
        from pptx import Presentation
        import io

        prs   = Presentation(io.BytesIO(file_bytes))
        texts = []
        slide_count = len(prs.slides)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

        full_text = "\n".join(texts)

        if not full_text:
            return AnalysisResponse(
                file_name=file_name, content_type="ppt", status="error",
                error_message="No text found in PowerPoint slides.",
                confidence=0.0
            )

        word_count = len(full_text.split())
        ai         = analyze_text_content(full_text, file_name)
        questions  = generate_practice_questions(full_text, file_name, count=5)
        plagiarism = check_plagiarism(full_text, file_name)

        return AnalysisResponse(
            file_name=file_name,
            content_type="ppt",
            summary=ai.get("summary"),
            tags=ai.get("tags", []),
            keywords=ai.get("keywords", []),
            difficulty=ai.get("difficulty"),
            language=ai.get("language", "English"),
            topics=ai.get("topics", []),
            suggested_year=ai.get("suggested_year"),
            confidence=ai.get("confidence", 0.8),
            quality_score=ai.get("quality_score"),
            quality_note=ai.get("quality_note"),
            recommendations=ai.get("recommendations", []),
            practice_questions=_to_question_models(questions),
            plagiarism_score=plagiarism.get("plagiarism_score"),
            plagiarism_note=plagiarism.get("plagiarism_note"),
            word_count=word_count,
            page_count=slide_count,
            status="success"
        )

    except ImportError:
        return AnalysisResponse(
            file_name=file_name, content_type="ppt", status="error",
            error_message="python-pptx not installed. Run: pip install python-pptx",
            confidence=0.0
        )


# ─── HELPER ──────────────────────────────────────────────────────────────────

def _to_question_models(raw: list) -> List[PracticeQuestion]:
    """Convert raw dicts from GPT to PracticeQuestion Pydantic models."""
    result = []
    for q in raw:
        try:
            result.append(PracticeQuestion(
                question=q.get("question", ""),
                options=q.get("options"),
                answer=q.get("answer"),
                question_type=q.get("question_type", "open"),
            ))
        except Exception:
            pass
    return result
